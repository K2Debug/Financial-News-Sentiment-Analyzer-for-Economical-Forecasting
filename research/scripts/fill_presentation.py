"""Copy and fill Year II Presentation II template with EF-02 project content."""
from __future__ import annotations

import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

SRC = Path(r"c:\Users\K2\Downloads\Year II Presantation II.pptx")
DST = Path(__file__).resolve().parents[1] / "Year II Presentation II.pptx"


def set_para_text(para, text: str) -> None:
  if para.runs:
    para.runs[0].text = text
    for run in para.runs[1:]:
      run.text = ""
  else:
    para.text = text


def set_shape_lines(shape, lines: list[str]) -> None:
  tf = shape.text_frame
  paras = list(tf.paragraphs)
  for i, line in enumerate(lines):
    if i < len(paras):
      set_para_text(paras[i], line)
    else:
      p = tf.add_paragraph()
      p.text = line
      if paras and paras[0].runs:
        r0 = paras[0].runs[0]
        if p.runs:
          p.runs[0].font.name = r0.font.name
          p.runs[0].font.size = r0.font.size
          p.runs[0].font.bold = r0.font.bold
  for j in range(len(lines), len(paras)):
    set_para_text(paras[j], "")


def add_table_row(table) -> None:
  tbl = table._tbl
  tbl.append(deepcopy(tbl.tr_lst[-1]))


def set_cell_text(cell, text: str) -> None:
  cell.text = text
  for para in cell.text_frame.paragraphs:
    for run in para.runs:
      if run.font.size is None:
        run.font.size = Pt(12)


def main() -> None:
  shutil.copy2(SRC, DST)
  prs = Presentation(str(DST))

  # Slide 1 — Title
  s1 = prs.slides[0]
  for shape in s1.shapes:
    if not shape.has_text_frame:
      continue
    name = shape.name
    if name == "Title 1":
      set_para_text(shape.text_frame.paragraphs[0], "EF-02: Financial News Sentiment Analyser for Economic Forecasting")
    elif name == "Subtitle 2" and "Supervisor" in shape.text_frame.text:
      set_shape_lines(
        shape,
        [
          "Supervisors:",
          "Dr. Ahmad",
          "Mr. Kamali",
        ],
      )
    elif name == "Subtitle 2" and "Group" in shape.text_frame.text:
      set_shape_lines(
        shape,
        [
          "Group Members:",
          "Kitula, Kitula-Masaganya Abdul (35552/T.2025)",
          "Ileka, Revocatus Orest (37058/T.2025)",
          "Lyellu, Sabrina Rashid (37463/T.2025)",
          "Ibrahim, Ghadhal Nkrumah (37254/T.2025)",
          "",
          "",
        ],
      )

  # Slide 2 — Methodology table
  methodology_rows = [
    (
      "1.",
      "To classify Tanzanian financial headlines by topic and economic sentiment",
      "Web scraping, cleaning, multi-model benchmarking, Groq LLM single-pass 3-in-1 classification",
      "Python, requests, BeautifulSoup, pandas, Groq API (Llama 3.1 8B)",
      "tz_headlines_labelled.csv (6,235 headlines; 82.8% relevant)",
    ),
    (
      "2.",
      "To compute monthly sentiment aggregates across 2022-2024",
      "Monthly grouping of relevant headlines; merge with NBS CPI and USD/TZS series",
      "pandas",
      "Visualization_Data.csv (36 monthly rows)",
    ),
    (
      "3.",
      "To correlate monthly sentiment with inflation and exchange rate",
      "Pearson correlation and significance testing on merged monthly dataset",
      "matplotlib, seaborn, scipy",
      "Correlation tables and dual-axis time series charts",
    ),
    (
      "4.",
      "To visualise sentiment trends against macroeconomic indicators",
      "Eight-section analysis notebook with relevance, sentiment, FX, and inflation charts",
      "matplotlib, seaborn",
      "06_visualisation.ipynb outputs",
    ),
    (
      "5.",
      "To evaluate NLP approaches for Tanzanian financial text",
      "Benchmark TextBlob, VADER, FinBERT, and LLMs on labelled sample",
      "TextBlob, VADER, transformers, Groq/OpenAI APIs",
      "benchmarking_results.csv (~93% LLM accuracy vs <67% for alternatives)",
    ),
  ]

  s2 = prs.slides[1]
  for shape in s2.shapes:
    if shape.has_table:
      table = shape.table
      while len(table.rows) < len(methodology_rows) + 1:
        add_table_row(table)
      for ri, row_data in enumerate(methodology_rows, start=1):
        for ci, value in enumerate(row_data):
          set_cell_text(table.cell(ri, ci), value)

  # Slide 3 — Users
  s3 = prs.slides[2]
  for shape in s3.shapes:
    if shape.name == "Content Placeholder 2":
      set_shape_lines(
        shape,
        [
          "Policymakers and economists monitoring currency pressure and inflation mood from local media",
          "Investors tracking sector trends (banking, trade, agriculture, forex) before official statistics",
          "Researchers studying African financial NLP and sentiment-macro relationships",
          "Data science students using the six-notebook pipeline as a reproducible case study",
        ],
      )

  # Slide 4 — Functionalities
  s4 = prs.slides[3]
  for shape in s4.shapes:
    if shape.name == "Content Placeholder 2":
      set_shape_lines(
        shape,
        [
          "Scrape Daily News and The Citizen business archives (2022-2024) with checkpoint resume",
          "Classify each headline for relevance, 11 economic categories, and sentiment by implied economic outcome",
          "Aggregate relevant headlines to monthly sentiment shares and dominant category",
          "Merge headline aggregates with NBS CPI inflation and USD/TZS exchange rate data",
          "Produce time series charts, Net Sentiment index, and Pearson correlation analysis",
          "Key finding: LLMs suitable for Tanzanian financial text; FX link suggestive (r = -0.29, p = 0.08) but not significant at 5%",
        ],
      )

  # Slide 5 — Demonstration
  s5 = prs.slides[4]
  for shape in s5.shapes:
    if shape.name == "Content Placeholder 2":
      set_shape_lines(
        shape,
        [
          "GitHub repository:",
          "https://github.com/K2Debug/Financial-News-Sentiment-Analyzer-for-Economic-Forecasting",
          "",
          "Run notebooks 01-06 in order (scraper to cleaning to classifier to consolidation to visualisation)",
        ],
      )

  prs.save(str(DST))
  print(f"Saved: {DST}")


if __name__ == "__main__":
  main()
